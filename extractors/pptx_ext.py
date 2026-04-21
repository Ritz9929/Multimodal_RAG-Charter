"""
PPTX Extractor
===============
Uses python-pptx for PowerPoint parsing.

Strategy:
  - Each slide → one "page" in ExtractionResult
  - Shape text (titles, body text, text boxes) → page text
  - Tables → markdown format
  - Images → extracted from picture shapes
  - Speaker notes → appended with [SPEAKER_NOTES] tag

Why include speaker notes?
  Presenters often put detailed explanations, data sources, and
  context in speaker notes. For RAG, this is high-value text that
  the slide itself may only reference via bullet points.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .base import BaseExtractor, ExtractionResult, ExtractedImage, ExtractedTable

logger = logging.getLogger(__name__)


class PPTXExtractor(BaseExtractor):
    """
    PPTX extraction using python-pptx.

    Strategy:
      - Each slide → one "page" in ExtractionResult
      - Shape text (titles, body text, text boxes) → page text
      - Tables → markdown format
      - Images → extracted from picture shapes
      - Speaker notes → appended with [SPEAKER_NOTES] tag
    """

    def __init__(self, output_dir: str = "./mock_s3_storage", **kwargs):
        super().__init__(output_dir)

    def _table_to_markdown(self, table) -> str:
        """Convert a python-pptx Table object to markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
        body_lines = []
        for row in rows[1:]:
            padded = row + [""] * (len(rows[0]) - len(row))
            body_lines.append("| " + " | ".join(padded[:len(rows[0])]) + " |")

        return "\n".join([header, separator] + body_lines)

    def _extract_shape_image(self, shape, slide_num: int, img_count: int,
                              doc_output_dir: Path) -> ExtractedImage | None:
        """Extract an image from a picture shape."""
        try:
            image = shape.image
            image_bytes = image.blob

            if len(image_bytes) < 1024:
                return None  # Skip tiny icons

            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"

            filename = f"slide{slide_num}_img{img_count}.{ext}"
            filepath = doc_output_dir / filename

            with open(filepath, "wb") as f:
                f.write(image_bytes)

            return ExtractedImage(
                filename=filename,
                filepath=str(filepath),
                page_number=slide_num,
                position_index=img_count,
                source_type="slide_image",
            )
        except Exception as e:
            logger.warning(f"  Skipping PPTX image on slide {slide_num}: {e}")
            return None

    def extract(self, file_path: str, source_doc: str = "") -> ExtractionResult:
        """Parse the PPTX and return text + images + tables per slide."""
        prs = Presentation(file_path)
        result = ExtractionResult(source_format="pptx")

        if not source_doc:
            source_doc = Path(file_path).name
        doc_output_dir = self._make_doc_output_dir(source_doc)

        logger.info(f"Opened PPTX: {file_path} ({len(prs.slides)} slides)")

        for slide_num, slide in enumerate(prs.slides):
            text_parts = []
            img_count = 0
            table_count = 0

            for shape in slide.shapes:
                # ── Text frames (titles, body text, text boxes) ──
                if shape.has_text_frame:
                    frame_text = shape.text_frame.text.strip()
                    if frame_text:
                        text_parts.append(frame_text)

                # ── Tables ──
                elif shape.has_table:
                    md = self._table_to_markdown(shape.table)
                    if md:
                        text_parts.append(md)
                        result.tables.append(ExtractedTable(
                            markdown=md,
                            page_number=slide_num,
                            position_index=table_count,
                            row_count=len(shape.table.rows),
                            col_count=len(shape.table.columns),
                        ))
                        table_count += 1

                # ── Images ──
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = self._extract_shape_image(
                        shape, slide_num, img_count, doc_output_dir
                    )
                    if img:
                        result.images.append(img)
                        img_count += 1

                # ── Group shapes (may contain nested text/images) ──
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            child_text = child.text_frame.text.strip()
                            if child_text:
                                text_parts.append(child_text)
                        elif child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            img = self._extract_shape_image(
                                child, slide_num, img_count, doc_output_dir
                            )
                            if img:
                                result.images.append(img)
                                img_count += 1

            # ── Speaker notes ──
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_parts.append(f"\n[SPEAKER_NOTES: {notes_text}]")

            # Store slide text
            result.page_texts[slide_num] = "\n".join(text_parts)

        result.metadata = {
            "slide_count": len(prs.slides),
            "tables_found": len(result.tables),
            "images_found": len(result.images),
        }

        logger.info(
            f"PPTX extraction complete — {len(prs.slides)} slides, "
            f"{len(result.tables)} tables, {len(result.images)} images"
        )
        return result
